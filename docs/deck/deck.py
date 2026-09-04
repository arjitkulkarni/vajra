# -*- coding: utf-8 -*-
"""VAJRA - SIH 2026 idea deck: template chrome, emblem, slides 1-3."""
import os
from pptx.util import Inches as In, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE as DASH
from kit import *

HERE = os.path.dirname(os.path.abspath(__file__))
SIH_LOGO = os.path.join(HERE, "logos", "sih.png")
DSU_LOGO = os.path.join(HERE, "logos", "dsu.png")
OUT = r"C:\Users\Arjit\Desktop\VAJRA - SIH2026 Idea Submission.pptx"

FOOT = "Idea Submission  \u00b7  Smart India Hackathon 2026 (Internal)  \u00b7  Dayananda Sagar University"


# ============================== chrome ==============================
def new_slide(prs, title, page, subhead=None, team=True):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = RGB(PAPER)

    if team:
        ov = circle(sl, M + 0.70, 0.66, 0.33, PAPER, PURPLE, 1.1)
        ov.width, ov.height = In(1.40), In(0.62)
        ov.left, ov.top = In(M), In(0.35)
        tb(sl, 0, 0, 0, 0, [P("The CodePool", SANS, 9.0, True, c=INK, al='c', ls=1.0)],
           anchor='m', shape_obj=ov)

    ts = fit_pt(title, 8.7, 26.0, SERIF, True, True, 0.0, minpt=15.0)
    tb(sl, 2.15, 0.30, 8.7, 0.5, [P(title, SERIF, ts, True, c=INK, al='c', caps=True, ls=1.0)])
    if subhead:
        ss = fit_pt(subhead, 8.7, 7.6, SANS, True, True, 1.2, minpt=6.0)
        tb(sl, 2.15, 0.84, 8.7, 0.2,
           [P(subhead, SANS, ss, True, c=BLUE, al='c', spc=1.2, caps=True, ls=1.0)])

    sl.shapes.add_picture(SIH_LOGO, In(SW - M - 1.66), In(0.30), In(1.66), In(0.783))
    hr(sl, M, RULE_Y, CW, RULE2, 1.0)

    rect(sl, 0, FOOT_Y, SW, SH - FOOT_Y, BLUE, None)
    tb(sl, 3.4, FOOT_Y + 0.155, 6.5, 0.22, [P(FOOT, SANS, 8.0, c=WHITE, al='c', ls=1.0)])
    tb(sl, SW - M - 1.4, FOOT_Y + 0.15, 1.4, 0.22,
       [P("%d" % page, SANS, 10.5, True, c=WHITE, al='r', ls=1.0)])
    return sl


def panel(sl, x, y, w, h, fill=PAPER2, line_=RULE2, lw=0.75, rad=0.03):
    return rect(sl, x, y, w, h, fill, line_, lw, rad)


# ============================== emblem ==============================
def vajra(sl, cx, cy, h, c=INK, acc=BLUE):
    u = h
    rect(sl, cx - u * 0.017, cy - u * 0.20, u * 0.034, u * 0.40, c, None)
    circle(sl, cx, cy, u * 0.070, PAPER, c, 1.2)
    circle(sl, cx, cy, u * 0.030, acc, None)
    for d in (-1, 1):
        rect(sl, cx - u * 0.085, cy + d * u * 0.205 - u * 0.014, u * 0.17, u * 0.028, c, None)
        rect(sl, cx - u * 0.055, cy + d * u * 0.135 - u * 0.010, u * 0.11, u * 0.020, c, None)
    for d in (-1, 1):
        base, tipy = cy + d * u * 0.225, cy + d * u * 0.470
        for s in (-1, 1):
            outer = qbez((cx + s * u * 0.048, base),
                         (cx + s * u * 0.215, base + d * u * 0.150),
                         (cx + s * u * 0.062, tipy), 16)
            inner = qbez((cx + s * u * 0.062, tipy),
                         (cx + s * u * 0.120, base + d * u * 0.155),
                         (cx + s * u * 0.014, base), 16)
            poly(sl, outer + inner, c, None)
        poly(sl, [(cx - u * 0.026, base), (cx, cy + d * u * 0.545), (cx + u * 0.026, base)],
             c, None)


def seal(sl, cx, cy, r):
    circle(sl, cx, cy, r, PAPER, BLUE, 1.6)
    circle(sl, cx, cy, r - 0.075, None, BLUEL, 0.9)
    for k in range(24):
        a = k * 15.0
        p0 = polar(cx, cy, r - 0.066, a)
        p1 = polar(cx, cy, r - (0.038 if k % 2 else 0.014), a)
        line(sl, p0[0], p0[1], p1[0], p1[1], BLUEM if k % 2 else BLUE, 0.8)
    shape(sl, MSO_SHAPE.HEXAGON, cx - r * 0.63, cy - r * 0.72, r * 1.26, r * 1.44,
          None, BLUEL, 0.9, rot=90)
    vajra(sl, cx, cy, r * 1.46)


# ============================== slide 1 ==============================
FIELDS = [("Problem Statement ID", "SIH26125", 0.44),
          ("Problem Statement Title",
           "Blockchain-Based Secure Platform for Identity, Access Control, and Digital Asset Management", 0.74),
          ("Theme", "Blockchain and Cyber Security", 0.44),
          ("PS Category", "Software", 0.44),
          ("Team ID", "SIH2026-070", 0.44),
          ("Team Name", "The CodePool", 0.44),
          ("Team Leader", "Arjit Kulkarni  \u00b7  USN ENG24CS0334", 0.44)]


def slide1(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    bg = sl.background.fill
    bg.solid()
    bg.fore_color.rgb = RGB(PAPER)

    sl.shapes.add_picture(DSU_LOGO, In(M), In(0.34), In(1.86), In(0.596))
    sl.shapes.add_picture(SIH_LOGO, In(SW - M - 1.80), In(0.28), In(1.80), In(0.849))
    tb(sl, 2.55, 0.30, 8.2, 0.46,
       [P("Smart India Hackathon 2026", SERIF, 27, True, c=INK, al='c', caps=True, ls=1.0)])
    tb(sl, 2.55, 0.78, 8.2, 0.34,
       [P("Title Page", SERIF, 19, True, c=BLUE, al='c', caps=True, ls=1.0)])
    hr(sl, M, RULE_Y + 0.06, CW, RULE2, 1.0)

    # ---- left: the required fields
    y = 1.62
    for lab, val, h in FIELDS:
        rect(sl, M + 0.04, y + 0.10, 0.085, 0.085, BLUE, None)
        tb(sl, M + 0.28, y, 6.35, h,
           [P([R(lab + "  \u2013  ", SANS, 11.5, b=True, c=INK),
               R(val, SANS, 11.5, b=False, c=INK2)], ls=1.22)])
        y += h

    for i, (t, sub) in enumerate([("\u20b90", "gas fees"), ("5", "live checks / request"),
                                  ("100%", "actions leave proof")]):
        bx = M + 0.05 + i * 2.20
        panel(sl, bx, 5.28, 2.02, 0.78, BLUEXL, BLUEL, 0.9, 0.06)
        tb(sl, bx, 5.38, 2.02, 0.32, [P(t, SANS, 19, True, c=BLUE, al='c', ls=1.0)])
        tb(sl, bx, 5.74, 2.02, 0.22,
           [P(sub, SANS, 7.4, True, c=INK2, al='c', spc=0.8, caps=True, ls=1.0)])

    tb(sl, M + 0.05, 6.24, 6.6, 0.44,
       [P("School of Engineering, Dayananda Sagar University", SANS, 9.6, True, c=INK, ls=1.2),
        P("Bengaluru \u2013 562112  \u00b7  Internal Hackathon 2026", SANS, 8.6, c=INK3, ls=1.2)])

    # ---- right: the product identity
    vr(sl, 7.30, 1.50, 5.10, RULE, 0.9)
    seal(sl, 10.22, 2.72, 1.12)
    tb(sl, 7.55, 3.96, 5.34, 0.62, [P("VAJRA", SERIF, 44, True, c=INK, al='c', spc=4.0, ls=1.0)])
    tb(sl, 7.55, 4.62, 5.34, 0.24,
       [P("Verifiable Authority & Zero-Trust Resource Architecture", SERIF, 11.6, i=True,
          c=INK3, al='c', ls=1.0)])
    hr(sl, 9.10, 4.98, 2.24, BLUE, 1.4)
    tb(sl, 7.55, 5.08, 5.34, 0.26,
       [P("\u201cTrust, Verified. Access, Controlled.\u201d", SERIF, 12.6, True, i=True,
          c=BLUE, al='c', ls=1.0)])
    for i, t in enumerate(["Identity", "Trust", "Decision", "Asset", "Proof"]):
        chip(sl, 7.55 + i * 1.08, 5.52, 1.02, 0.26, t, BLUEXL, BLUEL, BLUED, 7.0)
        if i < 4:
            tb(sl, 7.55 + i * 1.08 + 1.02, 5.53, 0.06, 0.22,
               [P("\u203a", SANS, 8.0, True, c=BLUEM, al='c', ls=1.0)])
    tb(sl, 7.55, 6.00, 5.34, 0.44,
       [P("Built and running today \u2014 117 unit tests, 87 end-to-end assertions,", SANS, 8.2,
          c=INK3, al='c', ls=1.22),
        P("on a laptop with no Docker and no accounts.", SANS, 8.2, c=INK3, al='c', ls=1.22)])

    rect(sl, 0, FOOT_Y, SW, SH - FOOT_Y, BLUE, None)
    tb(sl, 3.4, FOOT_Y + 0.155, 6.5, 0.22, [P(FOOT, SANS, 8.0, c=WHITE, al='c', ls=1.0)])
    tb(sl, SW - M - 1.4, FOOT_Y + 0.15, 1.4, 0.22,
       [P("1", SANS, 10.5, True, c=WHITE, al='r', ls=1.0)])
    return sl


# ============================== slide 2 ==============================
PROBLEMS = [(ic_key, "Passwords & shared credentials"),
            (ic_face_alert, "Deepfake & photo spoofing"),
            (ic_file_q, "No proof of custody"),
            (ic_audit_clock, "Week-long forensic audits"),
            (ic_db_lock, "Central biometric honeypots")]

INPUTS = [["Person", "Device", "Location", "Time", "Role"],
          ["Asset", "Action", "Risk", "Live proof", "Policy"]]

USPS = [(ic_trust, "Continuous Trust", "privileges shrink as trust drops"),
        (ic_asset, "Asset Passport", "owner \u00b7 lineage \u00b7 live trust score"),
        (ic_proof, "Proof-of-Action", "verifiable without trusting us"),
        (ic_shield, "Insider-Threat Response", "detect \u00b7 lock \u00b7 replay in one click"),
        (ic_audit_clock, "Time-Travel Audit", "what was true, at any moment")]


def slide2(prs):
    sl = new_slide(prs, "VAJRA \u2014 A Cryptographic Trust Layer", 2,
                   "Proposed solution  \u00b7  How it addresses the problem  \u00b7  Innovation and uniqueness")

    # ---- left : the problem
    x, w = M, 3.15
    eyebrow(sl, x, CY, w, "The problem today")
    ys, _ = stack([0.30] * 5, 1.78, 5.55, 0.30)
    for (fn, lab), yy in zip(PROBLEMS, ys):
        circle(sl, x + 0.20, yy + 0.15, 0.20, PAPER2, RULE2, 0.8)
        fn(sl, x + 0.20, yy + 0.15, 0.125, INK2)
        tb(sl, x + 0.52, yy - 0.02, w - 0.52, 0.36,
           [P(lab, SANS, 9.0, True, c=INK, ls=1.14)])

    # ---- centre : the trust firewall
    cx0, cw = 3.86, 5.20
    mid = cx0 + cw / 2
    eyebrow(sl, cx0, CY, cw, "The Trust Firewall \u2014 every request")
    cwid, cg = 0.98, 0.075
    for r, row in enumerate(INPUTS):
        for i, t in enumerate(row):
            chip(sl, cx0 + i * (cwid + cg), 1.80 + r * 0.30, cwid, 0.255, t,
                 BLUEXL, BLUEL, BLUED, 6.8)
    poly(sl, [(cx0, 2.46), (cx0 + cw, 2.46), (mid + 1.36, 2.86), (mid - 1.36, 2.86)],
         PAPER2, RULE2, 0.75)
    arrow_d(sl, mid, 2.88, 0.18, BLUE, 1.3)

    dh = 1.02
    shape(sl, MSO_SHAPE.DIAMOND, mid - 2.10, 3.08, 4.20, dh, BLUE, None)
    tb(sl, mid - 1.55, 3.34, 3.10, 0.24,
       [P("VAJRA decides", SANS, 11.0, True, c=WHITE, al='c', spc=0.6, caps=True, ls=1.0)])
    tb(sl, mid - 1.70, 3.60, 3.40, 0.22,
       [P("identity \u2227 RBAC \u2227 ABAC \u2227 risk \u2227 live proof", MONO, 7.2, c="CFE4F5",
          al='c', ls=1.0)])
    arrow_d(sl, mid, 4.12, 0.20, BLUE, 1.3)

    outs = [("Allow", BLUE, WHITE, None), ("Step-up", BLUEXL, BLUED, BLUEL),
            ("Deny", PAPER, DENY, DENY)]
    ow, og = 1.58, 0.23
    ox = mid - (3 * ow + 2 * og) / 2
    for i, (t, fill, tc, ln) in enumerate(outs):
        chip(sl, ox + i * (ow + og), 4.36, ow, 0.32, t, fill, ln, tc, 9.0, rad=0.4)
    arrow_d(sl, mid, 4.74, 0.20, BLUE, 1.3)

    # ledger strip
    bw, bg2, nb = 0.80, 0.30, 4
    bx0 = mid - (nb * bw + (nb - 1) * bg2) / 2
    for i in range(nb):
        bx = bx0 + i * (bw + bg2)
        rect(sl, bx, 5.02, bw, 0.42, PAPER, BLUED, 1.0, 0.10)
        tb(sl, bx, 5.11, bw, 0.24, [P("#%d" % (i + 1), MONO, 8.0, True, c=BLUED, al='c', ls=1.0)])
        if i < nb - 1:
            line(sl, bx + bw, 5.23, bx + bw + bg2, 5.23, BLUEM, 1.2)
    tb(sl, cx0, 5.54, cw, 0.22,
       [P("Anchored on chain \u00b7 hash-chained audit \u00b7 Proof-of-Action issued", SANS, 7.8,
          True, c=INK3, al='c', ls=1.0)])

    # ---- right : what's new
    rx, rw = 9.32, SW - M - 9.32
    eyebrow(sl, rx, CY, rw, "What makes it new")
    ys, _ = stack([0.56] * 5, 1.78, 5.55, 0.14)
    for i, ((fn, name, sub), yy) in enumerate(zip(USPS, ys)):
        panel(sl, rx, yy, rw, 0.56, PAPER2, RULE, 0.75, 0.05)
        circle(sl, rx + 0.30, yy + 0.28, 0.185, WHITE, BLUEL, 0.9)
        fn(sl, rx + 0.30, yy + 0.28, 0.115, BLUED)
        tb(sl, rx + 0.60, yy + 0.09, rw - 0.70, 0.42,
           [P(name, SANS, 8.6, True, c=INK, ls=1.10),
            P(sub, SANS, 7.4, c=INK3, ls=1.14)])
        tb(sl, rx + rw - 0.34, yy + 0.06, 0.30, 0.2,
           [P("0%d" % (i + 1), SANS, 7.6, True, c=BLUEM, al='r', ls=1.0)])

    # ---- bottom band
    by = 5.90
    panel(sl, M, by, CW, 0.90, BLUEXL, BLUEL, 0.9, 0.02)
    rect(sl, M, by, 0.07, 0.90, BLUE, None)
    tb(sl, M + 0.26, by + 0.13, 7.4, 0.2,
       [P("Every request is evaluated atomically, live", SANS, 7.6, True, c=BLUED,
          spc=1.2, caps=True, ls=1.0)])
    tb(sl, M + 0.26, by + 0.36, 7.4, 0.28,
       [P("Access = IdentityValid \u2227 RBAC \u2227 ABAC \u2227 (Risk < Threshold) \u2227 LiveProof",
          MONO, 12.6, True, c=INK, ls=1.0)])
    tb(sl, M + 0.26, by + 0.66, 7.4, 0.2,
       [P("Fail any one \u2014 and the action stops, with a named reason.", SANS, 7.6, i=True,
          c=INK3, ls=1.0)])
    vr(sl, 8.55, by + 0.14, 0.62, BLUEL, 1.2)
    tb(sl, 8.76, by + 0.14, SW - M - 8.90, 0.66,
       [P("\u201cVAJRA doesn\u2019t just control who can access an asset \u2014 it proves who accessed it, why they were allowed, and whether the asset can still be trusted.\u201d",
          SERIF, 10.2, i=True, c=INK2, ls=1.24)])
    return sl


# ============================== slide 3 ==============================
LAYERS = [(ic_users, "Users & Clients",
           "Next.js console \u00b7 admin dashboard \u00b7 live face capture \u00b7 en / hi / kn"),
          (ic_identity, "Identity Layer",
           "W3C DID minted in-browser \u00b7 credential issued \u00b7 nonce signed"),
          (ic_cloud, "Trust Gateway",
           "Fastify \u00b7 59 endpoints \u00b7 the only writer of record \u00b7 RLS reads"),
          (ic_decision, "Trust & Policy Engine",
           "RBAC \u2192 ABAC \u2192 trust gates \u2192 risk \u2192 approvals"),
          (ic_db_lock, "Data & Assets",
           "Neon Postgres, 27 tables \u00b7 AES-256-GCM \u00b7 content-addressed"),
          (ic_chain, "Blockchain Layer",
           "Fabric chaincode, 5 contracts \u00b7 hash-chained audit anchored")]

TECHS = [(ic_chain, "Blockchain", "Hyperledger Fabric v2.x\n5 Node.js contracts \u00b7 zero gas"),
         (ic_identity, "Biometrics", "AdaFace IR-50 + MiniFASNetV2\nin-browser, on device"),
         (ic_lock_ok, "Identity", "W3C DID \u00b7 Verifiable Credentials\nEd25519 \u00b7 EdDSA JWT"),
         (ic_trust, "AI risk engine", "Explainable scorer \u2014 every\npoint names its signal"),
         (ic_db_lock, "Database", "Neon serverless Postgres\n27 tables \u00b7 Drizzle ORM"),
         (ic_code, "Frontend", "Next.js 15 \u00b7 React 19\nTailwind v4 \u00b7 three languages"),
         (ic_cloud, "API layer", "Fastify + zod \u00b7 PostgREST\nPostgres row-level security"),
         (ic_server, "Storage", "IPFS / Pinata \u00b7 CIDv1\nonly the SHA-256 goes on chain")]

STEPS = [("Live identity", "face + liveness, on device"),
         ("Verifiable DID", "credential issued and checked"),
         ("Access request", "identity, action, asset, context"),
         ("Second live check", "re-verified at action time"),
         ("Trust + policy + risk", "ALLOW / STEP-UP / DENY"),
         ("Smart contract", "ownership + two-person rule"),
         ("Audit + proof", "anchored \u00b7 certificate issued")]


def slide3(prs):
    sl = new_slide(prs, "Technical Approach", 3,
                   "Technologies used  \u00b7  Methodology and process for implementation")

    LX, LW = M, 7.28
    eyebrow(sl, LX, CY, LW, "System architecture \u2014 six layers")
    y0, rh, pitch = 1.74, 0.56, 0.618
    vr(sl, LX + 0.14, y0 + 0.20, 5 * pitch, BLUEL, 1.4)
    for i, (fn, name, det) in enumerate(LAYERS):
        y = y0 + i * pitch
        hot = (i == 3)
        panel(sl, LX + 0.30, y, LW - 0.30, rh, BLUE if hot else PAPER2,
              None if hot else RULE, 0.75, 0.05)
        circle(sl, LX + 0.14, y + rh / 2, 0.14, BLUE if hot else WHITE, BLUE, 1.1)
        if not hot:
            tb(sl, LX, y + rh / 2 - 0.075, 0.28, 0.16,
               [P(str(i + 1), SANS, 7.4, True, c=BLUE, al='c', ls=1.0)])
        else:
            tb(sl, LX, y + rh / 2 - 0.075, 0.28, 0.16,
               [P(str(i + 1), SANS, 7.4, True, c=WHITE, al='c', ls=1.0)])
        fn(sl, LX + 0.58, y + rh / 2, 0.135, WHITE if hot else INK2)
        tb(sl, LX + 0.86, y + 0.10, 1.86, 0.24,
           [P(name, SANS, 9.2, True, c=WHITE if hot else INK, ls=1.0)])
        tb(sl, LX + 0.86, y + 0.32, LW - 1.20, 0.22,
           [P(det, SANS, 7.5, c="CFE4F5" if hot else INK3, ls=1.0)])
        if hot:
            for k, (t, fc, tc) in enumerate([("Allow", WHITE, BLUE), ("Step-up", None, WHITE),
                                             ("Deny", None, WHITE)]):
                chip(sl, LX + 4.95 + k * 0.78, y + 0.15, 0.72, 0.26, t, fc, WHITE, tc, 6.4, rad=0.4)
    ay = y0 + 6 * pitch - 0.06
    tb(sl, LX + 0.30, ay + 0.04, 1.30, 0.2,
       [P("Integrates with", SANS, 7.0, True, c=MUTED, spc=1.0, caps=True, ls=1.0)])
    for k, t in enumerate(["HR system", "ERP", "PLM / Design", "IAM / IdP", "Object store"]):
        chip(sl, LX + 1.62 + k * 1.14, ay, 1.08, 0.24, t, WHITE, RULE2, INK3, 6.6)

    # ---- technologies
    RX, RW = 8.02, SW - M - 8.02
    eyebrow(sl, RX, CY, RW, "Technology stack")
    tw, tg, th = (RW - 0.12) / 2, 0.12, 0.86
    for i, (fn, name, det) in enumerate(TECHS):
        tx = RX + (i % 2) * (tw + tg)
        ty = 1.74 + (i // 2) * (th + 0.10)
        panel(sl, tx, ty, tw, th, PAPER2, RULE, 0.75, 0.05)
        circle(sl, tx + 0.30, ty + 0.28, 0.185, WHITE, BLUEL, 0.9)
        fn(sl, tx + 0.30, ty + 0.28, 0.115, BLUED)
        tb(sl, tx + 0.58, ty + 0.15, tw - 0.68, 0.24,
           [P(name, SANS, 8.4, True, c=INK, ls=1.0)])
        tb(sl, tx + 0.16, ty + 0.46, tw - 0.30, 0.36,
           [P(det, SANS, 7.0, c=INK3, ls=1.16)])

    # ---- bottom band : 7 steps
    by = 5.86
    rect(sl, 0, by, SW, FOOT_Y - by, DEEP, None)
    tb(sl, M, by + 0.10, CW, 0.2,
       [P("Implementation workflow \u2014 seven steps, every one leaving evidence", SANS, 7.6,
          True, c=BLUEM, spc=1.2, caps=True, ls=1.0)])
    sw_, sg = 1.665, 0.135
    for i, (t, d) in enumerate(STEPS):
        sx = M + i * (sw_ + sg)
        circle(sl, sx + 0.135, by + 0.53, 0.135, BLUE, None)
        tb(sl, sx, by + 0.455, 0.27, 0.18,
           [P(str(i + 1), SANS, 7.6, True, c=WHITE, al='c', ls=1.0)])
        tb(sl, sx + 0.33, by + 0.40, sw_ - 0.36, 0.20,
           [P(t, SANS, 7.6, True, c=WHITE, ls=1.0)])
        tb(sl, sx + 0.33, by + 0.585, sw_ - 0.36, 0.30,
           [P(d, SANS, 6.6, c="9FB0BF", ls=1.14)])
        if i < 6:
            tb(sl, sx + sw_ + 0.005, by + 0.44, 0.13, 0.2,
               [P("\u203a", SANS, 11.0, True, c=BLUEM, al='c', ls=1.0)])
    tb(sl, M, by + 0.86, CW, 0.2,
       [P([R("The blockchain layer is the source of truth \u2014 if the off-chain cache is unavailable, the system fails closed, not open.    "),
           R("Judges, try this: ", c=BLUEM), R("fail the liveness check at step 4 and watch step 6 block instantly.", c=WHITE)],
          SANS, 7.4, i=True, c="9FB0BF", al='c', ls=1.0)])
    return sl
