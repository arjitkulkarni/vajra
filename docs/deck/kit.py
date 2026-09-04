# -*- coding: utf-8 -*-
"""VAJRA deck - drawing kit on the official SIH 2026 template chrome.

Palette and typography sampled from `final sih presentation.pdf`:
white ground, #006FC0 footer bar, #8063A1 team oval, black Times titles.
"""
from pptx.util import Inches as In, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE as DASH
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import (XL_CHART_TYPE, XL_LEGEND_POSITION, XL_MARKER_STYLE,
                             XL_TICK_LABEL_POSITION, XL_TICK_MARK)
import math

# ------------------------------- palette (template) -------------------------------
PAPER  = "FFFFFF"
PAPER2 = "F2F5F8"
PAPER3 = "E3E9EF"
WHITE  = "FFFFFF"
INK    = "111418"
INK2   = "333A42"
INK3   = "5C646E"
MUTED  = "949BA4"
RULE   = "DFE3E8"
RULE2  = "C3CAD2"
BLUE   = "006FC0"   # template footer bar - THE accent
BLUED  = "00548F"
BLUEM  = "4A9BD8"
BLUEL  = "A9D2EF"
BLUEXL = "E9F3FC"
PURPLE = "8063A1"   # template team oval
SLATE  = "44546A"   # SIH lockup grey
GREY   = "D8DEE5"
DEEP   = "14202B"
DEEP2  = "1E2C3A"
DEEP3  = "3D4E5E"
DENY   = "B4472E"
BRASS, BRASS2, BRASSL = BLUE, "6FB2E2", BLUEXL

SERIF = "Times New Roman"
SANS  = "Arial"
MONO  = "Consolas"

SW, SH = 13.3333, 7.5
M      = 0.44
CW     = SW - 2 * M
RULE_Y = 1.16
CY, CB = 1.32, 6.80
FOOT_Y = 6.94

GUT = 0.225
C3W = (CW - 2 * GUT) / 3
C3 = [M, M + C3W + GUT, M + 2 * (C3W + GUT)]


def RGB(h):
    return RGBColor.from_string(h)


# ------------------------------- primitives -------------------------------
def style(sh, fill=None, line=None, w=0.9, rad=None, dash=None):
    try:
        sh.shadow.inherit = False
    except Exception:
        pass
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = RGB(fill)
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = RGB(line)
        sh.line.width = Pt(w)
        if dash:
            sh.line.dash_style = dash
    if rad is not None:
        try:
            sh.adjustments[0] = rad
        except Exception:
            pass
    if sh.has_text_frame:
        tf = sh.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return sh


def rect(sl, x, y, w, h, fill=None, line=None, lw=0.9, rad=None, dash=None):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rad is not None else MSO_SHAPE.RECTANGLE
    sh = sl.shapes.add_shape(shp, In(x), In(y), In(w), In(h))
    return style(sh, fill, line, lw, rad, dash)


def shape(sl, kind, x, y, w, h, fill=None, line=None, lw=0.9, rad=None, rot=None):
    sh = sl.shapes.add_shape(kind, In(x), In(y), In(w), In(h))
    style(sh, fill, line, lw, rad)
    if rot:
        sh.rotation = rot
    return sh


def circle(sl, cx, cy, r, fill=None, line=None, lw=0.9):
    sh = sl.shapes.add_shape(MSO_SHAPE.OVAL, In(cx - r), In(cy - r), In(2 * r), In(2 * r))
    return style(sh, fill, line, lw)


def line(sl, x1, y1, x2, y2, c=RULE, w=0.75, dash=None):
    cn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, In(x1), In(y1), In(x2), In(y2))
    cn.line.color.rgb = RGB(c)
    cn.line.width = Pt(w)
    if dash:
        cn.line.dash_style = dash
    try:
        cn.shadow.inherit = False
    except Exception:
        pass
    return cn


def hr(sl, x, y, w, c=RULE, lw=0.75, dash=None):
    return line(sl, x, y, x + w, y, c, lw, dash)


def vr(sl, x, y, h, c=RULE, lw=0.75, dash=None):
    return line(sl, x, y, x, y + h, c, lw, dash)


def poly(sl, pts, fill=None, line_=None, lw=0.9, close=True):
    b = sl.shapes.build_freeform(Emu(int(pts[0][0] * 914400)), Emu(int(pts[0][1] * 914400)))
    b.add_line_segments([(Emu(int(p[0] * 914400)), Emu(int(p[1] * 914400))) for p in pts[1:]],
                        close=close)
    return style(b.convert_to_shape(), fill, line_, lw)


def qbez(p0, p1, p2, n=18):
    return [((1 - t) ** 2 * p0[0] + 2 * (1 - t) * t * p1[0] + t * t * p2[0],
             (1 - t) ** 2 * p0[1] + 2 * (1 - t) * t * p1[1] + t * t * p2[1])
            for t in (i / n for i in range(n + 1))]


def polar(cx, cy, r, a):
    return (cx + r * math.cos(math.radians(a)), cy + r * math.sin(math.radians(a)))


def arc_pts(cx, cy, r, a0, a1, n=40):
    if n <= 0:
        return [polar(cx, cy, r, a0)]
    return [polar(cx, cy, r, a0 + (a1 - a0) * i / n) for i in range(n + 1)]


def ring_arc(sl, cx, cy, r, a0, a1, c=INK, lw=1.0, n=40):
    pts = arc_pts(cx, cy, r, a0, a1, n)
    for i in range(len(pts) - 1):
        line(sl, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], c, lw)


# ------------------------------- text metrics -------------------------------
from PIL import ImageFont as _IF
import os as _os

_FDIR = _os.path.join(_os.environ.get("WINDIR", r"C:\Windows"), "Fonts")
_FF = {
    (SANS, 0, 0): "arial.ttf",    (SANS, 1, 0): "arialbd.ttf",
    (SANS, 0, 1): "ariali.ttf",   (SANS, 1, 1): "arialbi.ttf",
    (SERIF, 0, 0): "times.ttf",   (SERIF, 1, 0): "timesbd.ttf",
    (SERIF, 0, 1): "timesi.ttf",  (SERIF, 1, 1): "timesbi.ttf",
    (MONO, 0, 0): "consola.ttf",  (MONO, 1, 0): "consolab.ttf",
    (MONO, 0, 1): "consolai.ttf", (MONO, 1, 1): "consolaz.ttf",
}
_REF = 256.0
_CACHE = {}


def _font(f, bold=False, italic=False):
    key = (f, int(bool(bold)), int(bool(italic)))
    if key not in _CACHE:
        fn = _FF.get(key) or _FF.get((f, 0, 0)) or "arial.ttf"
        _CACHE[key] = _IF.truetype(_os.path.join(_FDIR, fn), int(_REF))
    return _CACHE[key]


def text_w(text, pt, f=SANS, bold=False, italic=False, spc=0.0):
    fo = _font(f, bold, italic)
    w = fo.getlength(text) / _REF * pt / 72.0
    return w + (len(text) * spc / 72.0 if spc else 0.0)


def wrap_lines(text, w_in, pt, f=SANS, bold=False, italic=False, spc=0.0, caps=False):
    src = text.upper() if caps else str(text)
    out = []
    for para in src.splitlines() or [""]:
        cur = ""
        for wd in para.split(" "):
            trial = wd if not cur else cur + " " + wd
            if text_w(trial, pt, f, bold, italic, spc) <= w_in or not cur:
                cur = trial
            else:
                out.append(cur)
                cur = wd
        out.append(cur)
    return out or [""]


def est_lines(text, w_in, pt, f=SANS, bold=False, spc=0.0, caps=False, italic=False,
              pad=1.045):
    return len(wrap_lines(text, w_in / pad, pt, f, bold, italic, spc, caps))


def est_h(text, w_in, pt, ls=1.16, f=SANS, bold=False, spc=0.0, caps=False, italic=False,
          pad=1.045):
    return est_lines(text, w_in, pt, f, bold, spc, caps, italic, pad) * pt / 72.0 * ls


def fit_pt(text, w_in, base_pt, f=SERIF, bold=True, caps=True, spc=0.0, minpt=11.0,
           italic=False):
    t = text.upper() if caps else text
    pt = base_pt
    while pt > minpt and text_w(t, pt, f, bold, italic, spc) > w_in:
        pt -= 0.25
    return pt


def stack(heights, top, bottom, gap_min=0.06):
    n = len(heights)
    if n == 0:
        return [], gap_min
    slack = bottom - top - sum(heights) - gap_min * (n - 1)
    g = gap_min + (max(0.0, slack) / (n - 1) if n > 1 else 0.0)
    ys, y = [], top
    for h in heights:
        ys.append(y)
        y += h + g
    return ys, g


# ------------------------------- text engine -------------------------------
_AL = {'l': PP_ALIGN.LEFT, 'c': PP_ALIGN.CENTER, 'r': PP_ALIGN.RIGHT, 'j': PP_ALIGN.JUSTIFY}
_AN = {'t': MSO_ANCHOR.TOP, 'm': MSO_ANCHOR.MIDDLE, 'b': MSO_ANCHOR.BOTTOM}
_NS = '{http://schemas.openxmlformats.org/drawingml/2006/main}'


def R(t, f=None, s=None, b=None, i=None, c=None, spc=None, caps=None):
    return dict(t=t, f=f, s=s, b=b, i=i, c=c, spc=spc, caps=caps)


def P(t, f=SANS, s=8.5, b=False, i=False, c=INK2, al='l', sa=0, sb=0, ls=1.16,
      spc=None, caps=False):
    return dict(t=t, f=f, s=s, b=b, i=i, c=c, al=al, sa=sa, sb=sb, ls=ls, spc=spc, caps=caps)


def _run(p, txt, f, s, b, i, c, spc, caps):
    r = p.add_run()
    r.text = txt.upper() if caps else txt
    fo = r.font
    fo.name = f
    fo.size = Pt(s)
    fo.bold = bool(b)
    fo.italic = bool(i)
    fo.color.rgb = RGB(c)
    rPr = fo._rPr
    rPr.set('lang', 'en-US')
    if spc:
        rPr.set('spc', str(int(spc * 100)))
    if rPr.find(_NS + 'cs') is None:
        anchor = rPr.find(_NS + 'ea')
        if anchor is None:
            anchor = rPr.find(_NS + 'latin')
        el = rPr.makeelement(_NS + 'cs', {'typeface': f})
        if anchor is not None:
            anchor.addnext(el)
        else:
            rPr.append(el)
    return r


def tb(sl, x, y, w, h, paras, anchor='t', wrap=True, shape_obj=None):
    if shape_obj is not None:
        tf = shape_obj.text_frame
    else:
        tf = sl.shapes.add_textbox(In(x), In(y), In(w), In(h)).text_frame
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = _AN[anchor]
    for idx, pa in enumerate(paras):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = _AL[pa['al']]
        if pa['sa']:
            p.space_after = Pt(pa['sa'])
        if pa['sb']:
            p.space_before = Pt(pa['sb'])
        if pa['ls']:
            p.line_spacing = pa['ls']
        t = pa['t']
        if isinstance(t, str):
            _run(p, t, pa['f'], pa['s'], pa['b'], pa['i'], pa['c'], pa['spc'], pa['caps'])
        else:
            for rd in t:
                _run(p, rd['t'], rd['f'] or pa['f'], rd['s'] or pa['s'],
                     pa['b'] if rd['b'] is None else rd['b'],
                     pa['i'] if rd['i'] is None else rd['i'],
                     rd['c'] or pa['c'],
                     pa['spc'] if rd['spc'] is None else rd['spc'],
                     pa['caps'] if rd['caps'] is None else rd['caps'])
    return tf


def label(sl, x, y, w, txt, c=MUTED, s=6.8, al='l', b=True, spc=1.2):
    return tb(sl, x, y, w, 0.17, [P(txt, SANS, s, b, c=c, al=al, spc=spc, caps=True, ls=1.0)])


def eyebrow(sl, x, y, w, txt, c=BLUE, s=10.0, rule=True):
    rect(sl, x, y + 0.012, 0.058, 0.16, c, None)
    tb(sl, x + 0.15, y, w - 0.15, 0.2,
       [P(txt, SANS, s, True, c=INK, spc=0.4, caps=True, ls=1.0)])
    if rule:
        hr(sl, x, y + 0.25, w, RULE2, 0.9)


def chip(sl, x, y, w, h, txt, fill=None, line_=RULE2, c=INK2, s=7.2, f=SANS, b=True,
         rad=0.5, spc=0.5, caps=True):
    sh = rect(sl, x, y, w, h, fill, line_, 0.75, rad)
    tb(sl, 0, 0, 0, 0, [P(txt, f, s, b, c=c, al='c', spc=spc, caps=caps, ls=1.0)],
       anchor='m', shape_obj=sh)
    return sh


def arrow_r(sl, x, y, ln, c=BLUE, lw=1.1, head=0.065):
    line(sl, x, y, x + ln - head, y, c, lw)
    poly(sl, [(x + ln - head * 1.9, y - head * 0.85), (x + ln, y),
              (x + ln - head * 1.9, y + head * 0.85)], c, None)


def arrow_d(sl, x, y, ln, c=BLUE, lw=1.1, head=0.065):
    line(sl, x, y, x, y + ln - head, c, lw)
    poly(sl, [(x - head * 0.85, y + ln - head * 1.9), (x, y + ln),
              (x + head * 0.85, y + ln - head * 1.9)], c, None)


def tick(sl, cx, cy, s=0.075, c=BLUE, lw=1.2):
    line(sl, cx - s, cy, cx - s * 0.25, cy + s * 0.72, c, lw)
    line(sl, cx - s * 0.25, cy + s * 0.72, cx + s * 1.02, cy - s * 0.78, c, lw)


def cross(sl, cx, cy, s=0.07, c=DENY, lw=1.2):
    line(sl, cx - s, cy - s, cx + s, cy + s, c, lw)
    line(sl, cx + s, cy - s, cx - s, cy + s, c, lw)


# ------------------------------- native charts -------------------------------
def _axis(ax, c=MUTED, s=7.0, grid=False, gridc=RULE, tick_labels=True):
    ax.has_major_gridlines = grid
    if grid:
        gl = ax.major_gridlines.format.line
        gl.color.rgb = RGB(gridc)
        gl.width = Pt(0.6)
    ax.format.line.color.rgb = RGB(RULE2)
    ax.format.line.width = Pt(0.75)
    ax.major_tick_mark = XL_TICK_MARK.NONE
    ax.minor_tick_mark = XL_TICK_MARK.NONE
    if tick_labels:
        ax.tick_labels.font.size = Pt(s)
        ax.tick_labels.font.name = SANS
        ax.tick_labels.font.color.rgb = RGB(c)
    else:
        ax.tick_label_position = XL_TICK_LABEL_POSITION.NONE


def line_chart(sl, x, y, w, h, cats, series, colors, ymin=0, ymax=100,
               label_series=0, legend=True, dash_idx=(), fsize=7.0):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    ch = sl.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, In(x), In(y), In(w), In(h), cd).chart
    ch.font.size = Pt(fsize)
    ch.font.name = SANS
    ch.font.color.rgb = RGB(INK3)
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(fsize)
    va, ca = ch.value_axis, ch.category_axis
    va.minimum_scale, va.maximum_scale = ymin, ymax
    _axis(va, MUTED, fsize, grid=True)
    _axis(ca, INK3, fsize)
    plot = ch.plots[0]
    for i, s in enumerate(plot.series):
        col = colors[i % len(colors)]
        s.format.line.color.rgb = RGB(col)
        s.format.line.width = Pt(1.1 if i in dash_idx else 2.4)
        s.smooth = False
        if i in dash_idx:
            s.format.line.dash_style = DASH.DASH
            s.marker.style = XL_MARKER_STYLE.NONE
        else:
            s.marker.style = XL_MARKER_STYLE.CIRCLE
            s.marker.size = 6
            s.marker.format.fill.solid()
            s.marker.format.fill.fore_color.rgb = RGB(col)
            s.marker.format.line.color.rgb = RGB(WHITE)
            s.marker.format.line.width = Pt(1.0)
    if label_series is not None:
        s = plot.series[label_series]
        for pt in s.points:
            dl = pt.data_label
            dl.has_text_frame = False
        s_dl = s.data_labels if hasattr(s, "data_labels") else None
        if s_dl is not None:
            s_dl.show_value = True
            s_dl.font.size = Pt(fsize - 0.3)
            s_dl.font.bold = True
            s_dl.font.name = SANS
            s_dl.font.color.rgb = RGB(INK)
            s_dl.number_format = '0'
            s_dl.number_format_is_linked = False
    return ch


def col_chart(sl, x, y, w, h, cats, series, colors, legend=False, gap=80,
              overlap=-12, fsize=7.0, ymax=None, labels=True, numfmt='0'):
    cd = CategoryChartData()
    cd.categories = cats
    for name, vals in series:
        cd.add_series(name, vals)
    ch = sl.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, In(x), In(y), In(w), In(h), cd).chart
    ch.font.size = Pt(fsize)
    ch.font.name = SANS
    ch.has_legend = legend
    if legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(fsize)
    plot = ch.plots[0]
    plot.gap_width = gap
    plot.overlap = overlap
    if labels:
        plot.has_data_labels = True
        dl = plot.data_labels
        dl.font.size = Pt(fsize)
        dl.font.bold = True
        dl.font.name = SANS
        dl.font.color.rgb = RGB(INK)
        dl.number_format = numfmt
        dl.number_format_is_linked = False
    va, ca = ch.value_axis, ch.category_axis
    if ymax:
        va.maximum_scale = ymax
    va.minimum_scale = 0
    _axis(va, MUTED, fsize, grid=True, tick_labels=False)
    _axis(ca, INK2, fsize)
    for i, s in enumerate(plot.series):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = RGB(colors[i % len(colors)])
        s.format.line.fill.background()
    return ch


def donut(sl, cx, cy, r, frac, c=BLUE, track=PAPER3, thick=0.30, ctext=None, s=13.0,
          sub=None):
    """ring gauge drawn from arc segments; frac 0..1"""
    n = 64
    for i in range(n):
        a0, a1 = -90 + 360.0 * i / n, -90 + 360.0 * (i + 1) / n
        col = c if (i / n) < frac else track
        p0, p1 = polar(cx, cy, r, a0), polar(cx, cy, r, a1)
        line(sl, p0[0], p0[1], p1[0], p1[1], col, thick * 72 / 1.0 * 0.0 + thick * 72)
    if ctext:
        tb(sl, cx - r, cy - 0.20, 2 * r, 0.26,
           [P(ctext, SANS, s, True, c=INK, al='c', ls=1.0)])
    if sub:
        tb(sl, cx - r, cy + 0.05, 2 * r, 0.20,
           [P(sub, SANS, 6.4, True, c=MUTED, al='c', spc=0.8, caps=True, ls=1.0)])


# ------------------------------- icons (monoline) -------------------------------
def ic_key(sl, cx, cy, s=0.16, c=INK):
    circle(sl, cx - s * 0.55, cy, s * 0.42, None, c, 1.0)
    circle(sl, cx - s * 0.55, cy, s * 0.15, c, None)
    line(sl, cx - s * 0.13, cy, cx + s * 1.05, cy, c, 1.0)
    line(sl, cx + s * 0.45, cy, cx + s * 0.45, cy + s * 0.42, c, 1.0)
    line(sl, cx + s * 0.85, cy, cx + s * 0.85, cy + s * 0.34, c, 1.0)


def ic_face_alert(sl, cx, cy, s=0.16, c=INK):
    for sx in (-1, 1):
        for sy in (-1, 1):
            line(sl, cx + sx * s * 1.02, cy + sy * s * 1.02,
                 cx + sx * s * 1.02, cy + sy * s * 0.58, c, 1.0)
            line(sl, cx + sx * s * 1.02, cy + sy * s * 1.02,
                 cx + sx * s * 0.58, cy + sy * s * 1.02, c, 1.0)
    circle(sl, cx, cy - s * 0.22, s * 0.34, None, c, 1.0)
    pts = qbez((cx - s * 0.62, cy + s * 0.72), (cx, cy + s * 0.02),
               (cx + s * 0.62, cy + s * 0.72), 14)
    for i in range(len(pts) - 1):
        line(sl, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], c, 1.0)
    line(sl, cx - s * 0.92, cy + s * 0.86, cx + s * 0.92, cy - s * 0.86, DENY, 1.5)


def ic_file_q(sl, cx, cy, s=0.16, c=INK):
    w, h = s * 0.92, s * 1.24
    poly(sl, [(cx - w / 2, cy - h / 2), (cx + w / 2 - s * 0.3, cy - h / 2),
              (cx + w / 2, cy - h / 2 + s * 0.3), (cx + w / 2, cy + h / 2),
              (cx - w / 2, cy + h / 2)], None, c, 1.0)
    line(sl, cx + w / 2 - s * 0.3, cy - h / 2, cx + w / 2 - s * 0.3, cy - h / 2 + s * 0.3, c, 0.8)
    line(sl, cx + w / 2 - s * 0.3, cy - h / 2 + s * 0.3, cx + w / 2, cy - h / 2 + s * 0.3, c, 0.8)
    tb(sl, cx - s * 0.4, cy - s * 0.30, s * 0.8, s * 0.72,
       [P("?", SANS, s * 42, True, c=c, al='c', ls=1.0)], anchor='m')


def ic_audit_clock(sl, cx, cy, s=0.16, c=INK):
    circle(sl, cx - s * 0.16, cy - s * 0.16, s * 0.62, None, c, 1.0)
    line(sl, cx - s * 0.16, cy - s * 0.16, cx - s * 0.16, cy - s * 0.54, c, 0.9)
    line(sl, cx - s * 0.16, cy - s * 0.16, cx + s * 0.16, cy - s * 0.06, c, 0.9)
    line(sl, cx + s * 0.30, cy + s * 0.30, cx + s * 0.86, cy + s * 0.86, c, 1.3)


def ic_db_lock(sl, cx, cy, s=0.16, c=INK):
    shape(sl, MSO_SHAPE.CAN, cx - s * 0.5, cy - s * 0.85, s * 1.0, s * 1.5, None, c, 1.0)
    rect(sl, cx + s * 0.24, cy + s * 0.16, s * 0.62, s * 0.5, PAPER, c, 0.9, 0.18)
    ring_arc(sl, cx + s * 0.55, cy + s * 0.18, s * 0.2, 180, 360, c, 0.9, 12)


def ic_identity(sl, cx, cy, s=0.16, c=INK):
    rect(sl, cx - s * 0.95, cy - s * 0.72, s * 1.9, s * 1.44, None, c, 1.0, 0.14)
    circle(sl, cx - s * 0.4, cy - s * 0.18, s * 0.26, None, c, 0.9)
    pts = qbez((cx - s * 0.82, cy + s * 0.42), (cx - s * 0.4, cy + s * 0.02),
               (cx + s * 0.02, cy + s * 0.42), 12)
    for i in range(len(pts) - 1):
        line(sl, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], c, 0.9)
    for k, ww in enumerate((0.52, 0.52, 0.32)):
        line(sl, cx + s * 0.24, cy - s * 0.24 + k * s * 0.26,
             cx + s * 0.24 + s * ww, cy - s * 0.24 + k * s * 0.26, c, 0.9)


def ic_trust(sl, cx, cy, s=0.16, c=INK):
    ring_arc(sl, cx, cy + s * 0.32, s * 0.86, 180, 360, c, 1.1, 26)
    line(sl, cx, cy + s * 0.32, cx + s * 0.5, cy - s * 0.22, BLUE, 1.4)
    circle(sl, cx, cy + s * 0.32, s * 0.1, c, None)
    for a in (200, 250, 290, 340):
        p0, p1 = polar(cx, cy + s * 0.32, s * 0.66, a), polar(cx, cy + s * 0.32, s * 0.78, a)
        line(sl, p0[0], p0[1], p1[0], p1[1], c, 0.8)


def ic_decision(sl, cx, cy, s=0.16, c=INK):
    shape(sl, MSO_SHAPE.DIAMOND, cx - s * 0.92, cy - s * 0.78, s * 1.84, s * 1.56, None, c, 1.0)
    tick(sl, cx - s * 0.02, cy - s * 0.02, s * 0.32, BLUE, 1.3)


def ic_asset(sl, cx, cy, s=0.16, c=INK):
    rect(sl, cx - s * 0.82, cy - s * 0.98, s * 1.64, s * 1.96, None, c, 1.0, 0.1)
    shape(sl, MSO_SHAPE.HEXAGON, cx - s * 0.42, cy - s * 0.72, s * 0.84, s * 0.74, None, BLUE, 1.1)
    for k, ww in enumerate((1.28, 1.28, 0.86)):
        line(sl, cx - s * 0.64, cy + s * 0.18 + k * s * 0.3,
             cx - s * 0.64 + s * ww, cy + s * 0.18 + k * s * 0.3, c, 0.8)


def ic_proof(sl, cx, cy, s=0.16, c=INK):
    circle(sl, cx, cy - s * 0.22, s * 0.72, None, c, 1.0)
    circle(sl, cx, cy - s * 0.22, s * 0.5, None, BLUE, 0.9)
    tick(sl, cx, cy - s * 0.24, s * 0.3, BLUE, 1.4)
    poly(sl, [(cx - s * 0.36, cy + s * 0.34), (cx - s * 0.36, cy + s * 1.06),
              (cx - s * 0.08, cy + s * 0.82), (cx + s * 0.2, cy + s * 1.06),
              (cx + s * 0.2, cy + s * 0.34)], None, c, 0.9)


def ic_shield(sl, cx, cy, s=0.16, c=INK):
    pts = [(cx, cy - s * 0.95), (cx + s * 0.78, cy - s * 0.6), (cx + s * 0.78, cy + s * 0.16)]
    pts += qbez((cx + s * 0.78, cy + s * 0.16), (cx + s * 0.6, cy + s * 0.8), (cx, cy + s * 1.0), 10)
    pts += qbez((cx, cy + s * 1.0), (cx - s * 0.6, cy + s * 0.8), (cx - s * 0.78, cy + s * 0.16), 10)
    pts += [(cx - s * 0.78, cy - s * 0.6)]
    poly(sl, pts, None, c, 1.0)


def ic_chain(sl, cx, cy, s=0.16, c=INK):
    rect(sl, cx - s * 0.95, cy - s * 0.3, s * 0.8, s * 0.6, None, c, 1.0, 0.5)
    rect(sl, cx + s * 0.15, cy - s * 0.3, s * 0.8, s * 0.6, None, c, 1.0, 0.5)
    line(sl, cx - s * 0.22, cy, cx + s * 0.22, cy, c, 1.2)


def ic_cloud(sl, cx, cy, s=0.16, c=INK):
    pts = [(cx - s * 0.76, cy + s * 0.44)]
    pts += arc_pts(cx - s * 0.40, cy + s * 0.08, s * 0.36, 180, 288, 12)
    pts += arc_pts(cx + s * 0.02, cy - s * 0.16, s * 0.46, 242, 322, 12)
    pts += arc_pts(cx + s * 0.50, cy + s * 0.12, s * 0.32, 262, 12, 12)
    pts += [(cx + s * 0.82, cy + s * 0.44)]
    poly(sl, pts, None, c, 1.0)


def ic_users(sl, cx, cy, s=0.16, c=INK):
    for dx in (-0.42, 0.42):
        circle(sl, cx + s * dx, cy - s * 0.30, s * 0.26, None, c, 0.9)
        pts = qbez((cx + s * dx - s * 0.44, cy + s * 0.52), (cx + s * dx, cy - s * 0.02),
                   (cx + s * dx + s * 0.44, cy + s * 0.52), 10)
        for i in range(len(pts) - 1):
            line(sl, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], c, 0.9)


def ic_server(sl, cx, cy, s=0.16, c=INK):
    for k in (-1, 0, 1):
        rect(sl, cx - s * 0.78, cy + k * s * 0.52 - s * 0.22, s * 1.56, s * 0.44, None, c, 0.95, 0.12)
        circle(sl, cx - s * 0.52, cy + k * s * 0.52, s * 0.07, c, None)


def ic_code(sl, cx, cy, s=0.16, c=INK):
    rect(sl, cx - s * 0.98, cy - s * 0.76, s * 1.96, s * 1.52, None, c, 1.0, 0.1)
    line(sl, cx - s * 0.98, cy - s * 0.38, cx + s * 0.98, cy - s * 0.38, c, 0.8)
    tb(sl, cx - s * 0.9, cy - s * 0.30, s * 1.8, s * 0.9,
       [P("</>", MONO, s * 36, True, c=c, al='c', ls=1.0)], anchor='m')


def ic_globe(sl, cx, cy, s=0.16, c=INK):
    circle(sl, cx, cy, s * 0.9, None, c, 1.0)
    line(sl, cx - s * 0.9, cy, cx + s * 0.9, cy, c, 0.8)
    for k in (0.34, 0.62):
        ring_arc(sl, cx, cy, s * 0.9, 90, 270, c, 0.0, 2)
    shape(sl, MSO_SHAPE.OVAL, cx - s * 0.36, cy - s * 0.9, s * 0.72, s * 1.8, None, c, 0.8)


def ic_lock_ok(sl, cx, cy, s=0.16, c=INK):
    rect(sl, cx - s * 0.68, cy - s * 0.12, s * 1.36, s * 0.98, None, c, 1.0, 0.14)
    ring_arc(sl, cx, cy - s * 0.14, s * 0.42, 180, 360, c, 1.0, 16)
    tick(sl, cx, cy + s * 0.36, s * 0.26, BLUE, 1.3)
